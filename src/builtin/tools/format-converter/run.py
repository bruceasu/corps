import argparse
import os
import subprocess
import sys
import tempfile
from pathlib import Path

def add_runtime_path() -> None:
    possible_dirs = []
    scripts_dir = os.getenv("CORPS_PYTHON_SCRIPTS_DIR")
    if scripts_dir:
        p = Path(scripts_dir).resolve()
        possible_dirs.append(p)
        possible_dirs.append(p / "_runtime")
    
    try:
        possible_dirs.append(Path(__file__).resolve().parents[3] / "_runtime")
    except Exception:
        pass

    for d in possible_dirs:
        if d.exists() and (d / "tool_runtime.py").exists():
            if str(d) not in sys.path:
                sys.path.insert(0, str(d))
            return

add_runtime_path()

try:
    from tool_runtime import emit_result, failure, success
except ImportError:
    def emit_result(res): print(res)
    def failure(tool, err): return {"ok": False, "error": err}
    def success(tool, out, data): return {"ok": True, "output": out, "data": data}

def run_pandoc(input_path: str, to_format: str) -> str:
    cmd = ["pandoc", input_path, "-t", "markdown" if to_format == "markdown" else "plain"]
    result = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="replace")
    if result.returncode == 0:
        return result.stdout
    else:
        raise Exception(f"Pandoc error: {result.stderr}")

def convert_docx(input_path: str) -> str:
    try:
        import docx
        doc = docx.Document(input_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except ImportError:
        return run_pandoc(input_path, "text")

def convert_xlsx(input_path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(input_path, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                lines.append("\t".join([str(c) if c is not None else "" for c in row]))
        return "\n".join(lines)
    except ImportError:
        return f"[Error: openpyxl not installed, cannot convert xlsx nicely]"

def convert_pptx(input_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(input_path)
        text_runs = []
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except ImportError:
        return run_pandoc(input_path, "text")

def convert_via_office_com(input_path: str, app_name: str) -> str:
    import win32com.client
    try:
        app = win32com.client.Dispatch(f"{app_name}.Application")
        if app_name == "Word":
            doc = app.Documents.Open(str(Path(input_path).absolute()))
            content = doc.Content.Text
            doc.Close()
            return content
        elif app_name == "Excel":
            wb = app.Workbooks.Open(str(Path(input_path).absolute()))
            content = ""
            for sheet in wb.Sheets:
                content += f"--- {sheet.Name} ---\n"
                # This is slow for large sheets, but okay for a general tool
                used_range = sheet.UsedRange
                vals = used_range.Value
                if vals:
                    for row in vals:
                        content += "\t".join([str(v) if v is not None else "" for v in row]) + "\n"
            wb.Close()
            return content
        elif app_name == "PowerPoint":
            pres = app.Presentations.Open(str(Path(input_path).absolute()), WithWindow=False)
            text_runs = []
            for slide in pres.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            text_runs.append(shape.TextFrame.TextRange.Text)
            pres.Close()
            return "\n".join(text_runs)
    except Exception as e:
        raise e
    finally:
        try: app.Quit()
        except: pass

def main():
    parser = argparse.ArgumentParser(description="Convert files to Markdown/Text.")
    parser.add_argument("--input", required=True)
    parser.add_argument("--to", default="markdown")
    args = parser.parse_args()

    input_path = args.input
    to_format = args.to.lower()

    if not os.path.exists(input_path):
        emit_result(failure("format-converter", f"File not found: {input_path}"))
        return

    ext = Path(input_path).suffix.lower()
    content = ""
    
    try:
        if ext == ".docx":
            content = convert_docx(input_path)
        elif ext == ".xlsx":
            content = convert_xlsx(input_path)
        elif ext == ".pptx":
            content = convert_pptx(input_path)
        elif ext in [".doc", ".xls", ".ppt"]:
            # These old formats really need Office COM
            try:
                app_map = {".doc": "Word", ".xls": "Excel", ".ppt": "PowerPoint"}
                content = convert_via_office_com(input_path, app_map[ext])
            except Exception:
                # Fallback to pandoc if possible
                content = run_pandoc(input_path, to_format)
        elif ext == ".pdf":
            content = run_pandoc(input_path, to_format)
        else:
            # Fallback for other formats
            content = run_pandoc(input_path, to_format)

        emit_result(success("format-converter", content, {
            "input": input_path,
            "to": to_format,
            "size": len(content)
        }))

    except Exception as e:
        emit_result(failure("format-converter", f"Conversion failed: {str(e)}"))

if __name__ == "__main__":
    main()
